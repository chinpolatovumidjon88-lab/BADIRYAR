"""
Jan Bodriyar g'oyalari asosida virtual reallik va axloqiy muammolar tahlili
PowerPoint prezentatsiyasi yaratish skripti
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Rang palitrasi (zamonaviy, akademik)
DARK_BLUE = RGBColor(0x0A, 0x2A, 0x43)      # asosiy fon
ACCENT = RGBColor(0xE8, 0xB5, 0x4A)         # oltin urg'u
LIGHT = RGBColor(0xF5, 0xF5, 0xF0)          # och matn
GRAY = RGBColor(0xB0, 0xB8, 0xC1)           # ikkilamchi matn
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_background(slide, color=DARK_BLUE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=LIGHT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_slide_header(slide, title, number, total):
    add_accent_bar(slide, Inches(0.6), Inches(0.55))
    add_text(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.7),
             title, size=28, bold=True, color=WHITE)
    add_text(slide, Inches(11.8), Inches(0.5), Inches(1.3), Inches(0.4),
             f"{number} / {total}", size=12, color=GRAY, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.25),
                                  Inches(12.13), Emu(12700))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


# ============================================================
# SLIDE 1 — Sarlavha
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

# Dekorativ shakllar
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(10.5), Inches(-1.5),
                                Inches(4.5), Inches(4.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT
circle.line.fill.background()

small = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(-1), Inches(5.5),
                               Inches(3), Inches(3))
small.fill.solid()
small.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x5C)
small.line.fill.background()

add_accent_bar(slide, Inches(1.0), Inches(2.4), width=Inches(0.12), height=Inches(2.2))

add_text(slide, Inches(1.3), Inches(2.3), Inches(10), Inches(0.5),
         "FALSAFIY TAHLIL", size=14, bold=True, color=ACCENT)

add_text(slide, Inches(1.3), Inches(2.8), Inches(10.5), Inches(2.5),
         "Jan Bodriyar g'oyalari asosida\nvirtual reallik va axloqiy muammolar tahlili",
         size=36, bold=True, color=WHITE)

add_text(slide, Inches(1.3), Inches(5.3), Inches(10), Inches(0.5),
         "Simulyakr • Giperreallik • Axloq • Raqamli davr",
         size=16, color=GRAY)

add_text(slide, Inches(1.3), Inches(6.6), Inches(10), Inches(0.4),
         "Prezentatsiya  |  2026", size=11, color=GRAY)


# ============================================================
# SLIDE 2 — Reja
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT)
add_slide_header_dark = lambda s, t, n, tot: None  # placeholder

# Light variant header
add_accent_bar(slide, Inches(0.6), Inches(0.55))
add_text(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.7),
         "Reja", size=28, bold=True, color=DARK_BLUE)
add_text(slide, Inches(11.8), Inches(0.5), Inches(1.3), Inches(0.4),
         "2 / 10", size=12, color=GRAY, align=PP_ALIGN.RIGHT)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.25), Inches(12.13), Emu(12700))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

reja = [
    "Jan Bodriyar haqida qisqacha ma'lumot",
    "Asosiy tushunchalar: simulyakr va simulyatsiya",
    "Giperreallik nazariyasi",
    "Tasvirning to'rt bosqichi",
    "Virtual reallikning falsafiy tabiati",
    "Axloqiy muammolar: haqiqat va yolg'on chegarasi",
    "Zamonaviy raqamli olamga ta'siri",
    "Xulosa va munozara",
]

# 2 ustun ko'rinishida
left_items = reja[:4]
right_items = reja[4:]

for i, item in enumerate(left_items):
    y = Inches(1.9 + i * 1.05)
    num = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.8), y, Inches(0.7), Inches(0.7))
    num.fill.solid()
    num.fill.fore_color.rgb = DARK_BLUE
    num.line.fill.background()
    add_text(slide, Inches(0.8), y + Inches(0.13), Inches(0.7), Inches(0.5),
             f"{i+1:02d}", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.7), y + Inches(0.18), Inches(5), Inches(0.6),
             item, size=15, color=DARK_TEXT)

for i, item in enumerate(right_items):
    y = Inches(1.9 + i * 1.05)
    num = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(7.0), y, Inches(0.7), Inches(0.7))
    num.fill.solid()
    num.fill.fore_color.rgb = ACCENT
    num.line.fill.background()
    add_text(slide, Inches(7.0), y + Inches(0.13), Inches(0.7), Inches(0.5),
             f"{i+5:02d}", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.9), y + Inches(0.18), Inches(5), Inches(0.6),
             item, size=15, color=DARK_TEXT)


# ============================================================
# SLIDE 3 — Bodriyar haqida
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_slide_header(slide, "Jan Bodriyar (Jean Baudrillard) kim?", 3, 10)

# Chap blok — biografiya
add_text(slide, Inches(0.85), Inches(1.6), Inches(6), Inches(0.5),
         "Qisqacha ma'lumot", size=14, bold=True, color=ACCENT)

bio = [
    "1929–2007 — Frantsuz faylasufi va sotsiologi",
    "Postmodern falsafa va madaniyat tanqidchisi",
    "Sorbonna universitetida sotsiologiya bo'yicha tahsil olgan",
    "Iste'molchilik jamiyati va ommaviy axborot tahlili bilan shug'ullangan",
    "G'oyalari kino, san'at va texnologiyaga katta ta'sir ko'rsatgan",
]
add_bullets(slide, Inches(0.85), Inches(2.1), Inches(6.5), Inches(4),
            bio, size=15, color=LIGHT)

# O'ng blok — asosiy asarlar
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(7.7), Inches(1.6),
                              Inches(5.0), Inches(5.4))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x5C)
card.line.fill.background()

add_text(slide, Inches(8.0), Inches(1.85), Inches(4.5), Inches(0.5),
         "Asosiy asarlari", size=16, bold=True, color=ACCENT)

asarlar = [
    ("1970", "Iste'molchilik jamiyati"),
    ("1981", "Simulyakr va simulyatsiya"),
    ("1991", "Fors ko'rfazi urushi bo'lmagan"),
    ("1995", "Mukammal jinoyat"),
    ("2000", "Yovuzlikning shaffofligi"),
]
for i, (year, name) in enumerate(asarlar):
    y = Inches(2.5 + i * 0.8)
    add_text(slide, Inches(8.0), y, Inches(0.9), Inches(0.4),
             year, size=14, bold=True, color=ACCENT)
    add_text(slide, Inches(9.0), y, Inches(3.7), Inches(0.4),
             name, size=13, color=LIGHT)


# ============================================================
# SLIDE 4 — Simulyakr va simulyatsiya
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT)

add_accent_bar(slide, Inches(0.6), Inches(0.55))
add_text(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.7),
         "Asosiy tushunchalar: Simulyakr va Simulyatsiya",
         size=26, bold=True, color=DARK_BLUE)
add_text(slide, Inches(11.8), Inches(0.5), Inches(1.3), Inches(0.4),
         "4 / 10", size=12, color=GRAY, align=PP_ALIGN.RIGHT)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.25), Inches(12.13), Emu(12700))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Ikki katta katak
card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.85), Inches(1.7),
                               Inches(5.8), Inches(5.3))
card1.fill.solid()
card1.fill.fore_color.rgb = DARK_BLUE
card1.line.fill.background()

add_text(slide, Inches(1.15), Inches(1.95), Inches(5.5), Inches(0.5),
         "SIMULYAKR", size=14, bold=True, color=ACCENT)
add_text(slide, Inches(1.15), Inches(2.4), Inches(5.5), Inches(0.6),
         "Nusxasi bo'lmagan nusxa", size=20, bold=True, color=WHITE)
add_text(slide, Inches(1.15), Inches(3.2), Inches(5.5), Inches(3.5),
         "Bodriyar fikricha, simulyakr — bu asl narsani aks "
         "ettirmaydigan, balki o'zi mustaqil ravishda mavjud "
         "bo'lgan tasvir yoki belgidir.\n\n"
         "U haqiqatdan ajralib chiqib, o'zining alohida \"haqiqati\"ni "
         "yaratadi. Belgilar endi narsalarni emas, boshqa belgilarni "
         "ifodalaydi.",
         size=13, color=LIGHT)

card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(6.85), Inches(1.7),
                               Inches(5.8), Inches(5.3))
card2.fill.solid()
card2.fill.fore_color.rgb = ACCENT
card2.line.fill.background()

add_text(slide, Inches(7.15), Inches(1.95), Inches(5.5), Inches(0.5),
         "SIMULYATSIYA", size=14, bold=True, color=DARK_BLUE)
add_text(slide, Inches(7.15), Inches(2.4), Inches(5.5), Inches(0.6),
         "Haqiqatni almashtirish jarayoni", size=20, bold=True, color=DARK_BLUE)
add_text(slide, Inches(7.15), Inches(3.2), Inches(5.5), Inches(3.5),
         "Simulyatsiya — bu shunday holatki, unda asl bilan nusxa, "
         "haqiqat bilan tasavvur o'rtasidagi farq yo'qoladi.\n\n"
         "Natijada inson o'zi yashayotgan dunyoning haqiqiymi yoki "
         "yo'qligini aniqlay olmay qoladi. Bu — postmodern shart-"
         "sharoitning markaziy holatidir.",
         size=13, color=DARK_TEXT)


# ============================================================
# SLIDE 5 — Giperreallik
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_slide_header(slide, "Giperreallik nazariyasi", 5, 10)

add_text(slide, Inches(0.85), Inches(1.55), Inches(12), Inches(0.6),
         "\"Giperreallik — bu haqiqatdan ham haqiqiyroq bo'lib tuyuladigan reallik.\"",
         size=18, bold=True, color=ACCENT)

add_text(slide, Inches(0.85), Inches(2.35), Inches(12), Inches(1.2),
         "Bodriyar fikricha, zamonaviy jamiyatda ommaviy axborot vositalari, reklama va "
         "raqamli texnologiyalar shunday tasvirlar yaratadiki, ular asl voqelikdan "
         "kuchliroq, jonliroq va ishonchliroq ko'rinadi. Inson endi haqiqatga emas, "
         "uning tasviriga ishonadi.",
         size=14, color=LIGHT)

# 3 ta misol
examples = [
    ("DISNEYLAND", "Sun'iy yaratilgan dunyo bo'lib, Amerika hayotini \"haqiqiyroq\" qilib ko'rsatadi."),
    ("REALITY-SHOULAR", "Tahrir qilingan, sahnalashtirilgan voqealar tomoshabin uchun haqiqat sifatida taqdim etiladi."),
    ("IJTIMOIY TARMOQLAR", "Foydalanuvchilar o'z hayotining \"yaxshilangan\" versiyasini ko'rsatadi va o'zlari ham unga ishonadi."),
]

for i, (title, desc) in enumerate(examples):
    x = Inches(0.85 + i * 4.05)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(4.0), Inches(3.9), Inches(2.9))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x5C)
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.5)

    add_text(slide, x + Inches(0.25), Inches(4.2), Inches(3.6), Inches(0.5),
             title, size=14, bold=True, color=ACCENT)
    add_text(slide, x + Inches(0.25), Inches(4.75), Inches(3.6), Inches(2.0),
             desc, size=12, color=LIGHT)


# ============================================================
# SLIDE 6 — Tasvirning to'rt bosqichi
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT)

add_accent_bar(slide, Inches(0.6), Inches(0.55))
add_text(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.7),
         "Tasvirning to'rt bosqichi (Bodriyar)",
         size=26, bold=True, color=DARK_BLUE)
add_text(slide, Inches(11.8), Inches(0.5), Inches(1.3), Inches(0.4),
         "6 / 10", size=12, color=GRAY, align=PP_ALIGN.RIGHT)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.25), Inches(12.13), Emu(12700))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

stages = [
    ("I", "Aks ettirish",
     "Tasvir chuqur haqiqatni\naks ettiradi. Belgi va\nvoqelik mos keladi.",
     RGBColor(0x2E, 0x7D, 0x32)),
    ("II", "Buzib ko'rsatish",
     "Tasvir haqiqatni\nyashiradi va buzadi.\nBelgi haqiqatni soxtalashtiradi.",
     RGBColor(0xF9, 0xA8, 0x25)),
    ("III", "Yashirish",
     "Tasvir haqiqatning\nyo'qligini niqoblaydi.\nU \"bormi\" degan illyuziya yaratadi.",
     RGBColor(0xE6, 0x4A, 0x19)),
    ("IV", "Sof simulyakr",
     "Tasvir haqiqatga\numuman bog'liq emas.\nU o'zining sof simulyakridir.",
     RGBColor(0xC6, 0x28, 0x28)),
]

for i, (num, title, desc, color) in enumerate(stages):
    x = Inches(0.7 + i * 3.1)
    # Yuqori rangli bosh qism
    head = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(1.7), Inches(2.9), Inches(1.3))
    head.fill.solid()
    head.fill.fore_color.rgb = color
    head.line.fill.background()
    add_text(slide, x, Inches(1.85), Inches(2.9), Inches(0.6),
             num, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.5), Inches(2.9), Inches(0.5),
             title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Pastki tavsif
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(3.1), Inches(2.9), Inches(3.5))
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = GRAY
    body.line.width = Pt(0.75)
    add_text(slide, x + Inches(0.25), Inches(3.35), Inches(2.4), Inches(3.0),
             desc, size=13, color=DARK_TEXT)

# Strelka effekti (pastda)
add_text(slide, Inches(0.85), Inches(6.85), Inches(12), Inches(0.4),
         "Haqiqatdan uzoqlashish yo'nalishi  →  Sof simulyatsiyaga o'tish",
         size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — Virtual reallikning falsafiy tabiati
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_slide_header(slide, "Virtual reallikning falsafiy tabiati", 7, 10)

add_text(slide, Inches(0.85), Inches(1.55), Inches(12), Inches(0.6),
         "Bodriyar nuqtai nazaridan virtual reallik (VR)",
         size=18, bold=True, color=ACCENT)

points = [
    ("Haqiqatning oxiri",
     "VR — bu Bodriyar bashorat qilgan giperreallikning eng yuqori ko'rinishi. "
     "Foydalanuvchi tanasidan, makondan va hatto vaqtdan ajralib, sof simulyatsiyaga kiradi."),
    ("Tana va ong dixotomiyasi",
     "VR ongni jismoniy tanadan \"ajratadi\". Bu insonning o'zi haqidagi tushunchasini, "
     "identifikatsiyasini va mavjudligini qayta ko'rib chiqishga majbur qiladi."),
    ("Ramziy almashinuvning yo'qolishi",
     "An'anaviy madaniyatda belgilar haqiqiy munosabatlarni ifodalagan. "
     "VRda esa belgilar faqat boshqa belgilarga ishora qiladi — yopiq tizim paydo bo'ladi."),
    ("Yangi \"haqiqat\" turi",
     "VR — yolg'on emas, lekin haqiqat ham emas. U — uchinchi turdagi mavjudlik: "
     "texnologik vositachilik orqali yaratilgan giper-haqiqat."),
]

for i, (title, desc) in enumerate(points):
    y = Inches(2.3 + i * 1.18)
    # Raqam doirasi
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.85), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_text(slide, Inches(0.85), y + Inches(0.07), Inches(0.55), Inches(0.4),
             str(i+1), size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1.6), y - Inches(0.02), Inches(11), Inches(0.4),
             title, size=15, bold=True, color=ACCENT)
    add_text(slide, Inches(1.6), y + Inches(0.35), Inches(11), Inches(0.85),
             desc, size=12, color=LIGHT)


# ============================================================
# SLIDE 8 — Axloqiy muammolar
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT)

add_accent_bar(slide, Inches(0.6), Inches(0.55))
add_text(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.7),
         "Axloqiy muammolar: haqiqat va yolg'on chegarasi",
         size=24, bold=True, color=DARK_BLUE)
add_text(slide, Inches(11.8), Inches(0.5), Inches(1.3), Inches(0.4),
         "8 / 10", size=12, color=GRAY, align=PP_ALIGN.RIGHT)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(1.25), Inches(12.13), Emu(12700))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

problems = [
    ("Haqiqat tushunchasining yemirilishi",
     "Agar simulyatsiya haqiqatdan farq qilmasa, yolg'on va rost o'rtasidagi "
     "axloqiy chegara qaerda? Inson nimaga ishonishi kerak?"),
    ("Mas'uliyat masalasi",
     "Virtual makonda qilingan harakatlar (zo'ravonlik, aldash) uchun kim javob beradi? "
     "Ular axloqiy jihatdan haqiqiymi?"),
    ("Identifikatsiya va shaxs",
     "Avatarlar va raqamli ego — bu \"men\"ning haqiqiy ifodasimi yoki yana bir simulyakrmi?"),
    ("Empatiya va begonalashuv",
     "Ekran orqali muloqot insoniy hamdardlikni kamaytiradi. Ekranda azob ko'rgan boshqalar — "
     "endi shunchaki tasvir."),
    ("Manipulyatsiya xavfi",
     "Hokimiyat va korporatsiyalar giperreallik orqali jamoatchilik fikrini boshqarishi mumkin "
     "(masalan, deepfake, soxta yangiliklar)."),
    ("Erkin iroda muammosi",
     "Algoritmlar va simulyatsiyalar bizning tanlovimizni shakllantirsa, biz haqiqatan ham erkinmizmi?"),
]

# 2x3 grid
for i, (title, desc) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.7 + col * 4.15)
    y = Inches(1.55 + row * 2.7)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(3.95), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.25)

    # Chap chiziq
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(0.12), Inches(2.5))
    side.fill.solid()
    side.fill.fore_color.rgb = DARK_BLUE
    side.line.fill.background()

    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.5), Inches(0.6),
             title, size=14, bold=True, color=DARK_BLUE)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.5), Inches(1.6),
             desc, size=11, color=DARK_TEXT)


# ============================================================
# SLIDE 9 — Zamonaviy raqamli olamga ta'siri
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_slide_header(slide, "Zamonaviy raqamli olamga ta'siri", 9, 10)

add_text(slide, Inches(0.85), Inches(1.55), Inches(12), Inches(0.6),
         "Bodriyar g'oyalari bugun qanchalik dolzarb?",
         size=16, bold=True, color=ACCENT)

# Chap — ijtimoiy tarmoqlar va deepfake
card_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.85), Inches(2.2), Inches(5.9), Inches(4.7))
card_l.fill.solid()
card_l.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x5C)
card_l.line.fill.background()

add_text(slide, Inches(1.1), Inches(2.4), Inches(5.5), Inches(0.5),
         "RAQAMLI HAYOT", size=14, bold=True, color=ACCENT)

digital = [
    "Ijtimoiy tarmoqlardagi \"ideal hayot\" — sof simulyakr",
    "Deepfake va sun'iy intellekt yaratgan kontent — haqiqatni shubha ostiga oladi",
    "Metaverse — Bodriyar bashoratining moddiylashishi",
    "Onlayn shaxsiyatlar (avatar, profil) — yangi \"men\" turlari",
    "Algoritmik puzyrlar — har bir foydalanuvchi o'z giperrealligida",
]
add_bullets(slide, Inches(1.1), Inches(2.95), Inches(5.5), Inches(3.8),
            digital, size=12, color=LIGHT)

# O'ng — siyosat va madaniyat
card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.95), Inches(2.2), Inches(5.9), Inches(4.7))
card_r.fill.solid()
card_r.fill.fore_color.rgb = ACCENT
card_r.line.fill.background()

add_text(slide, Inches(7.2), Inches(2.4), Inches(5.5), Inches(0.5),
         "MADANIYAT VA SIYOSAT", size=14, bold=True, color=DARK_BLUE)

culture = [
    "Soxta yangiliklar (fake news) — haqiqat siyosiy qurolga aylandi",
    "Reklama va brendlar — narsani emas, tasvirni sotadi",
    "Urush va inqirozlar — ekran orqali \"shou\"ga aylanadi",
    "Madaniyat takror ishlab chiqarish (remake, remix) ustidan qurilgan",
    "Tarixiy haqiqat — turli versiyalarda \"qayta yoziladi\"",
]
add_bullets(slide, Inches(7.2), Inches(2.95), Inches(5.5), Inches(3.8),
            culture, size=12, color=DARK_TEXT)

# Pastdagi sitata
quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.85), Inches(7.0), Inches(12), Inches(0.4))
quote.fill.solid()
quote.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x5C)
quote.line.fill.background()


# ============================================================
# SLIDE 10 — Xulosa
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

# Dekorativ
deco = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(-2), Inches(-2), Inches(5), Inches(5))
deco.fill.solid()
deco.fill.fore_color.rgb = RGBColor(0x14, 0x3D, 0x5C)
deco.line.fill.background()

deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(11), Inches(5), Inches(4), Inches(4))
deco2.fill.solid()
deco2.fill.fore_color.rgb = ACCENT
deco2.line.fill.background()

add_accent_bar(slide, Inches(1.0), Inches(0.8), height=Inches(0.8))
add_text(slide, Inches(1.3), Inches(0.85), Inches(11), Inches(0.7),
         "Xulosa", size=36, bold=True, color=WHITE)

# Asosiy xulosa matni
add_text(slide, Inches(1.3), Inches(2.0), Inches(11), Inches(0.5),
         "Jan Bodriyar bizga muhim ogohlantirish qoldirdi:",
         size=16, bold=True, color=ACCENT)

conclusions = [
    "Virtual reallik shunchaki texnologiya emas — bu falsafiy va axloqiy hodisa.",
    "Haqiqat va simulyatsiya o'rtasidagi chegara yo'qolib borayotgan davrda inson "
    "tanqidiy fikrlashni saqlab qolishi kerak.",
    "Axloq endi nafaqat haqiqiy dunyoda, balki raqamli makonda ham qayta tiklanishi lozim.",
    "Bodriyar g'oyalari — bu pessimizm emas, balki ongli yashashga chaqiriqdir.",
]

for i, c in enumerate(conclusions):
    y = Inches(2.7 + i * 0.85)
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(1.3), y + Inches(0.1),
                                  Inches(0.25), Inches(0.25))
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT
    icon.line.fill.background()
    add_text(slide, Inches(1.75), y, Inches(10.5), Inches(0.8),
             c, size=14, color=LIGHT)

# Yakuniy sitata
add_text(slide, Inches(1.3), Inches(6.5), Inches(11), Inches(0.5),
         "\"Endi haqiqat yo'q — faqat uning simulyatsiyasi bor.\"  — J. Bodriyar",
         size=14, bold=True, color=ACCENT)

add_text(slide, Inches(1.3), Inches(7.0), Inches(11), Inches(0.4),
         "E'tiboringiz uchun rahmat!", size=18, bold=True, color=WHITE)


# ============================================================
output = "/projects/sandbox/Bodriyar_Virtual_Reallik.pptx"
prs.save(output)
print(f"Tayyor! Fayl saqlandi: {output}")
print(f"Jami slaydlar: {len(prs.slides)}")
