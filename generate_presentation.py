"""
Jan Bodriyar g'oyalari asosida virtual reallik va axloqiy muammolar tahlili
15 betlik ilmiy taqdimot
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BLUE = RGBColor(0x1D, 0x3B, 0x5C)
ACCENT_RED = RGBColor(0xC7, 0x3E, 0x1D)
ACCENT_GOLD = RGBColor(0xE9, 0xC4, 0x6A)
LIGHT_BG = RGBColor(0xF1, 0xFA, 0xEE)
TEXT_DARK = RGBColor(0x26, 0x46, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)


def add_background(slide, color):
    """Add solid color background"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_header_bar(slide, title_text, slide_num=None, total=15):
    """Add decorative header bar"""
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Gold accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), prs.slide_width, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_GOLD
    stripe.line.fill.background()

    # Title text
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(11.5), Inches(0.6))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'

    # Slide number
    if slide_num:
        nb = slide.shapes.add_textbox(Inches(12.0), Inches(0.25), Inches(1.2), Inches(0.4))
        ntf = nb.text_frame
        ntf.margin_left = 0
        np_ = ntf.paragraphs[0]
        np_.text = f"{slide_num} / {total}"
        np_.font.size = Pt(13)
        np_.font.color.rgb = ACCENT_GOLD
        np_.font.bold = True
        np_.alignment = PP_ALIGN.RIGHT


def add_footer(slide):
    """Bottom footer"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), prs.slide_width, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(12.5), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Jan Bodriyar g'oyalari asosida virtual reallik va axloqiy muammolar tahlili"
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_GOLD
    p.font.italic = True


def add_text_block(slide, left, top, width, height, text, font_size=14,
                   bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    return tb


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    bullet_color=ACCENT_RED, text_color=TEXT_DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▶  " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.space_after = Pt(8)
        p.font.name = 'Calibri'
    return tb


def add_decorated_box(slide, left, top, width, height, title, content,
                      title_color=DARK_BLUE, fill_color=LIGHT_BG):
    """Box with title and content"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = title_color
    box.line.width = Pt(1.5)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                        width, Inches(0.5))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = title_color
    title_bar.line.fill.background()

    title_tb = title_bar.text_frame
    title_tb.margin_left = Inches(0.15)
    p = title_tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    ctb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.55),
                                    width - Inches(0.3), height - Inches(0.6))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    if isinstance(content, list):
        for i, item in enumerate(content):
            p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)
    else:
        p = ctf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK


# ===================== SLIDE 1: TITLE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

# Decorative shapes
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = ACCENT_RED
circle1.line.fill.background()
circle1.element.getparent().remove(circle1.element)
slide.shapes._spTree.insert(2, circle1.element)

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = ACCENT_GOLD
circle2.line.fill.background()

# Top label
add_text_block(slide, Inches(1), Inches(0.6), Inches(11.3), Inches(0.5),
               "ILMIY-FALSAFIY TADQIQOT  |  MAGISTRLIK DISSERTATSIYASI",
               font_size=14, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.2), Inches(2.3), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_GOLD
line.line.fill.background()

# Main title
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0),
               "JAN BODRIYAR G'OYALARI ASOSIDA",
               font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_block(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5),
               "VIRTUAL REALLIK VA AXLOQIY",
               font_size=42, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

add_text_block(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(1.0),
               "MUAMMOLAR TAHLILI",
               font_size=42, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Subtitle
add_text_block(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(0.5),
               "«Simulyakr va simulyatsiya nazariyasi raqamli asrning ontologik-axloqiy parametrlarida»",
               font_size=15, italic=True, color=LIGHT_BG, align=PP_ALIGN.CENTER)

# Bottom info
info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(2.5), Inches(5.6), Inches(8.3), Inches(1.2))
info_box.fill.solid()
info_box.fill.fore_color.rgb = WHITE
info_box.line.color.rgb = ACCENT_GOLD
info_box.line.width = Pt(2)

add_text_block(slide, Inches(2.7), Inches(5.7), Inches(7.9), Inches(0.4),
               "Falsafa fanlari bo'yicha magistrlik dissertatsiyasi",
               font_size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text_block(slide, Inches(2.7), Inches(6.1), Inches(7.9), Inches(0.4),
               "Mualif: _______________________   |   Ilmiy rahbar: _______________________",
               font_size=12, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text_block(slide, Inches(2.7), Inches(6.45), Inches(7.9), Inches(0.4),
               "Toshkent — 2026",
               font_size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ===================== SLIDE 2: KIRISH - DOLZARBLIK =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "KIRISH: TADQIQOTNING DOLZARBLIGI", slide_num=2)
add_footer(slide)

# Intro paragraph
intro_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.4), Inches(1.15), Inches(12.5), Inches(1.4))
intro_box.fill.solid()
intro_box.fill.fore_color.rgb = WHITE
intro_box.line.color.rgb = DARK_BLUE
intro_box.line.width = Pt(1.5)

add_text_block(slide, Inches(0.6), Inches(1.25), Inches(12.1), Inches(1.2),
               "XXI asrning birinchi choragi insoniyat tarixida texnologik inqilob davri sifatida tavsiflanadi. "
               "Sun'iy intellekt, virtual reallik (VR), kengaytirilgan reallik (AR) va metakoinot kabi raqamli muhitlar "
               "insonning ontologik mavjudligi va axloqiy mo'ljallarini tubdan o'zgartirmoqda. Aynan shu sharoitda "
               "fransuz faylasufi Jan Bodriyar (1929–2007) ning simulyakr va giperreallik nazariyasi alohida ahamiyat "
               "kasb etadi.",
               font_size=13, color=TEXT_DARK)

# Three pillars of relevance
pillars = [
    ("FALSAFIY", "Klassik ontologiya va\ngnoseologiya tushunchalari\nraqamli asr sharoitida\nqayta ko'rib chiqilmoqda",
     ACCENT_RED),
    ("IJTIMOIY", "Ijtimoiy tarmoqlar,\nfeyk-news va deepfake\ntexnologiyalari haqiqat\ntushunchasini yemiradi",
     TEAL),
    ("AXLOQIY", "Virtual zo'ravonlik,\nraqamli identitet va\ninsoniy begonalashuv\nyangi axloqiy dilemmalar",
     ACCENT_GOLD),
]

for i, (title, content, color) in enumerate(pillars):
    left = Inches(0.4 + i * 4.25)
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, Inches(2.85), Inches(4.05), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    add_text_block(slide, left, Inches(2.95), Inches(4.05), Inches(0.4),
                   title + " DOLZARBLIK", font_size=15, bold=True,
                   color=WHITE, align=PP_ALIGN.CENTER)

    # Content body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, Inches(3.45), Inches(4.05), Inches(2.3))
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = color
    body.line.width = Pt(1.5)
    add_text_block(slide, left + Inches(0.15), Inches(3.6), Inches(3.75), Inches(2.0),
                   content, font_size=12.5, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# Bottom highlight
hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.0))
hl.fill.solid()
hl.fill.fore_color.rgb = DARK_BLUE
hl.line.fill.background()
add_text_block(slide, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.4),
               "STATISTIK ASOS:",
               font_size=13, bold=True, color=ACCENT_GOLD)
add_text_block(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
               "2025-yilda VR foydalanuvchilari soni 850 mln dan oshdi, metakoinot iqtisodiyoti 800 mlrd $ ga yetdi. "
               "Bu — Bodriyar bashorat qilgan «giperreal davr»ning empirik tasdig'idir.",
               font_size=12, italic=True, color=WHITE)


# ===================== SLIDE 3: MAQSAD VA VAZIFALAR =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "TADQIQOT MAQSADI, OB'EKTI VA VAZIFALARI", slide_num=3)
add_footer(slide)

# Object and Subject
add_decorated_box(slide, Inches(0.4), Inches(1.1), Inches(6.2), Inches(1.7),
                  "TADQIQOT OB'EKTI",
                  "Jan Bodriyarning simulyakr, simulyatsiya va giperreallik konsepsiyalari hamda zamonaviy raqamli-virtual muhitda namoyon bo'ladigan axloqiy-ontologik muammolar majmui.",
                  title_color=DARK_BLUE)

add_decorated_box(slide, Inches(6.7), Inches(1.1), Inches(6.2), Inches(1.7),
                  "TADQIQOT PREDMETI",
                  "Bodriyar falsafasidagi virtual reallik fenomenining ontologik maqomi va uning insoniy ekzistensiyaga, axloqiy ongga ta'sirining nazariy-falsafiy parametrlari.",
                  title_color=ACCENT_RED)

# Goal
goal_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.4), Inches(2.95), Inches(12.5), Inches(1.0))
goal_box.fill.solid()
goal_box.fill.fore_color.rgb = ACCENT_GOLD
goal_box.line.color.rgb = DARK_BLUE
goal_box.line.width = Pt(2)

add_text_block(slide, Inches(0.6), Inches(3.05), Inches(12.1), Inches(0.4),
               "TADQIQOT MAQSADI:",
               font_size=14, bold=True, color=DARK_BLUE)
add_text_block(slide, Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.6),
               "Jan Bodriyarning simulyatsiya nazariyasi prizmasidan virtual reallikning ontologik tabiati va uning axloqiy "
               "implikatsiyalarini kompleks falsafiy tahlil qilish, hamda muammoni hal etish bo'yicha nazariy yondashuvlar ishlab chiqish.",
               font_size=12.5, color=DARK_BLUE)

# Tasks
add_text_block(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.4),
               "TADQIQOTNING ASOSIY VAZIFALARI:",
               font_size=15, bold=True, color=DARK_BLUE)

tasks_left = [
    "Bodriyar falsafasining intellektual genezisi va nazariy manbalarini aniqlash",
    "Simulyakrning to'rt bosqich nazariyasini sistematik tahlil qilish",
    "Virtual reallikning ontologik maqomini falsafiy aniqlash",
]
tasks_right = [
    "Raqamli giperreallik sharoitidagi axloqiy dilemmalarni aniqlash",
    "Virtual zo'ravonlikning insoniy ongga ta'sirini o'rganish",
    "Bodriyar g'oyalariga tanqidiy yondashuv ishlab chiqish",
]

for i, t in enumerate(tasks_left):
    add_text_block(slide, Inches(0.5), Inches(4.55 + i*0.55), Inches(6.0), Inches(0.5),
                   f"{i+1}.  {t}", font_size=11.5, color=TEXT_DARK)

for i, t in enumerate(tasks_right):
    add_text_block(slide, Inches(6.8), Inches(4.55 + i*0.55), Inches(6.0), Inches(0.5),
                   f"{i+4}.  {t}", font_size=11.5, color=TEXT_DARK)

# Methodology
method = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.7))
method.fill.solid()
method.fill.fore_color.rgb = TEAL
method.line.fill.background()
add_text_block(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
               "METODOLOGIYA:  germenevtik-fenomenologik tahlil  •  qiyosiy-tarixiy metod  •  diskurs-tahlil  •  poststrukturalistik dekonstruksiya",
               font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ===================== SLIDE 4: ILMIY YANGILIK =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "TADQIQOTNING ILMIY YANGILIGI", slide_num=4)
add_footer(slide)

# Top description
add_text_block(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.6),
               "Mazkur tadqiqot Bodriyar falsafasini O'zbekiston falsafiy fikrida birinchi marta XXI asr raqamli "
               "muhiti kontekstida kompleks tahlil qiladi va quyidagi yangiliklarni taklif etadi:",
               font_size=13, italic=True, color=TEXT_DARK)

# 6 novelty cards in 2 rows
novelties = [
    ("01", "KONSEPTUAL YANGILIK",
     "Bodriyarning «giperreal-axloq» tushunchasi ilk bor sistemali shakllantirildi va falsafiy diskursga kiritildi.",
     ACCENT_RED),
    ("02", "METODOLOGIK YANGILIK",
     "Simulyatsiya nazariyasi va virtual reallik fenomenologiyasini birlashtirgan integrativ metod taklif qilindi.",
     TEAL),
    ("03", "TIPOLOGIK YANGILIK",
     "Virtual zo'ravonlikning yangi tasnifi: simulyativ, immersiv, algoritmik va onto-axloqiy turlari ajratildi.",
     ACCENT_GOLD),
    ("04", "EMPIRIK YANGILIK",
     "Metaverse va deepfake fenomenlari Bodriyar nazariyasi orqali falsafiy talqin qilindi.",
     DARK_BLUE),
    ("05", "TANQIDIY YANGILIK",
     "Bodriyarning texno-pessimistik qarashlari muqobil sharqona axloq prizmasidan qayta baholandi.",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("06", "AMALIY YANGILIK",
     "Raqamli muhitda axloqiy ta'lim uchun konseptual model ishlab chiqildi va tavsiya etildi.",
     RGBColor(0xE7, 0x6F, 0x51)),
]

for i, (num, title, content, color) in enumerate(novelties):
    row = i // 3
    col = i % 3
    left = Inches(0.4 + col * 4.25)
    top = Inches(1.75 + row * 2.45)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, Inches(4.05), Inches(2.25))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     left + Inches(0.2), top + Inches(0.15),
                                     Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_text_block(slide, left + Inches(0.2), top + Inches(0.27),
                   Inches(0.7), Inches(0.5),
                   num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title
    add_text_block(slide, left + Inches(1.0), top + Inches(0.25),
                   Inches(2.95), Inches(0.5),
                   title, font_size=12.5, bold=True, color=color)

    # Content
    add_text_block(slide, left + Inches(0.25), top + Inches(0.95),
                   Inches(3.65), Inches(1.25),
                   content, font_size=11, color=TEXT_DARK)


# ===================== SLIDE 5: NAZARIY MANBALAR =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "I BOB. BODRIYAR FALSAFASINING NAZARIY MANBALARI", slide_num=5)
add_footer(slide)

# Center figure
center = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(5.4), Inches(3.0), Inches(2.5), Inches(2.5))
center.fill.solid()
center.fill.fore_color.rgb = ACCENT_RED
center.line.color.rgb = ACCENT_GOLD
center.line.width = Pt(3)

add_text_block(slide, Inches(5.4), Inches(3.4), Inches(2.5), Inches(0.5),
               "JAN", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_block(slide, Inches(5.4), Inches(3.85), Inches(2.5), Inches(0.5),
               "BODRIYAR", font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_block(slide, Inches(5.4), Inches(4.35), Inches(2.5), Inches(0.5),
               "1929–2007", font_size=12, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Surrounding theoretical sources
sources = [
    ("MARKSIZM", "Iste'mol jamiyati,\nbegonalashuv tahlili\n(K. Marks, A. Lefebvr)", Inches(0.4), Inches(1.2), TEAL),
    ("STRUKTURALIZM", "Belgilar tizimi va\nsemiotika nazariyasi\n(F. de Sossyur, R. Bart)", Inches(9.5), Inches(1.2), DARK_BLUE),
    ("POSTSTRUKTURALIZM", "Dekonstruksiya va\nhokimiyat-bilim\n(J. Derrida, M. Fuko)", Inches(0.4), Inches(3.3), ACCENT_RED),
    ("PSIXOANALIZ", "Istakning tabiati\nva ramziy tartib\n(Z. Freyd, J. Lakan)", Inches(9.5), Inches(3.3), RGBColor(0x8E, 0x44, 0xAD)),
    ("FENOMENOLOGIYA", "Tajriba va idrok\nfalsafasi\n(E. Husserl, M. Merlo-Ponti)", Inches(0.4), Inches(5.4), ACCENT_GOLD),
    ("MEDIA NAZARIYALARI", "Media-soblidir,\ntexnomadaniyat\n(M. Maklyuen)", Inches(9.5), Inches(5.4), RGBColor(0xE7, 0x6F, 0x51)),
]

for name, desc, left, top, color in sources:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(3.4), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Title strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, Inches(3.4), Inches(0.45))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()
    add_text_block(slide, left, top + Inches(0.07), Inches(3.4), Inches(0.4),
                   name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_block(slide, left + Inches(0.15), top + Inches(0.55),
                   Inches(3.1), Inches(1.0),
                   desc, font_size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Connecting lines (decorative)
    # Skipped for cleanliness - boxes themselves are enough

# Bottom note
note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(3.5), Inches(7.2), Inches(6.3), Inches(0.0))


# ===================== SLIDE 6: GIPERREALLIK =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "1.2. GIPERREALLIK VA HAQIQAT TANAZZULI", slide_num=6)
add_footer(slide)

# Quote at top
quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.95))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = DARK_BLUE
quote_box.line.fill.background()

add_text_block(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.4),
               '"Giperreallik — bu haqiqatdan ham haqiqiyroq bo\'lgan modellashtirilgan voqelikdir."',
               font_size=14, bold=True, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
add_text_block(slide, Inches(0.7), Inches(1.65), Inches(12.0), Inches(0.4),
               "— J. Bodriyar, «Simulyakr va simulyatsiya», 1981",
               font_size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# Definition
add_decorated_box(slide, Inches(0.4), Inches(2.25), Inches(6.2), Inches(2.0),
                  "GIPERREALLIK — TUSHUNCHANING MOHIYATI",
                  ["Belgilar va modellar haqiqatni almashtirgan holat",
                   "Asl nusxa va nusxa o'rtasidagi farq yo'qoladi",
                   "Simulyakrlar avtonom tarzda yashay boshlaydi",
                   "Ontologik referent yo'qoladi — belgilar mustaqil reallik bo'ladi"],
                  title_color=ACCENT_RED)

# Examples
add_decorated_box(slide, Inches(6.7), Inches(2.25), Inches(6.2), Inches(2.0),
                  "ZAMONAVIY MISOLLAR",
                  ["Disneylend (Bodriyar klassik misoli)",
                   "Instagram va TikTok filtri haqiqati",
                   "Deepfake va sun'iy intellekt kontenti",
                   "Metaverse va VR-olamlar (Meta, VRChat)",
                   "AI-influencerlar (Lil Miquela, Aitana)"],
                  title_color=TEAL)

# Three phases of reality decay
add_text_block(slide, Inches(0.4), Inches(4.45), Inches(12.5), Inches(0.4),
               "HAQIQAT TANAZZULINING UCH BOSQICHI (Bodriyar bo'yicha):",
               font_size=14, bold=True, color=DARK_BLUE)

phases = [
    ("MODERNITET", "Sanoat asri\nMashina takrorlash", "Belgilar haqiqatni\naks ettiradi", ACCENT_GOLD),
    ("POSTMODERNITET", "Kommunikatsiya asri\nElektron media", "Belgilar haqiqatni\nbuzadi", ACCENT_RED),
    ("GIPERMODERNITET", "Raqamli asr\nVR/AR/Metaverse", "Belgilar HAQIQATNI\nALMASHTIRADI", DARK_BLUE),
]

for i, (era, char, role, color) in enumerate(phases):
    left = Inches(0.4 + i * 4.25)
    top = Inches(4.95)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(4.05), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, Inches(4.05), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_block(slide, left, top + Inches(0.1), Inches(4.05), Inches(0.4),
                   era, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_block(slide, left + Inches(0.2), top + Inches(0.6),
                   Inches(3.65), Inches(0.6),
                   char, font_size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_block(slide, left + Inches(0.2), top + Inches(1.2),
                   Inches(3.65), Inches(0.7),
                   role, font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Arrow between phases
    if i < 2:
        arrow_x = left + Inches(4.05)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         arrow_x - Inches(0.1), top + Inches(0.85),
                                         Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_RED
        arrow.line.fill.background()


# ===================== SLIDE 7: SIMULYAKRNING 4 BOSQICHI (CHART) =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "SIMULYAKRNING TO'RT BOSQICH NAZARIYASI", slide_num=7)
add_footer(slide)

# Insert chart
slide.shapes.add_picture('charts/chart1_simulacra_stages.png',
                          Inches(0.3), Inches(1.0),
                          width=Inches(8.5), height=Inches(5.2))

# Side explanation
side_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(9.0), Inches(1.0), Inches(4.0), Inches(5.2))
side_box.fill.solid()
side_box.fill.fore_color.rgb = WHITE
side_box.line.color.rgb = DARK_BLUE
side_box.line.width = Pt(2)

add_text_block(slide, Inches(9.15), Inches(1.1), Inches(3.7), Inches(0.4),
               "NAZARIY TUSHUNTIRISH",
               font_size=14, bold=True, color=DARK_BLUE)

stages_text = [
    ("I bosqich:", "Sodda, samimiy mashinaviy nusxa"),
    ("II bosqich:", "Sanoat asri — yolg'on aks ettirish, mafkura"),
    ("III bosqich:", "Aks ettiriladigan haqiqat o'zi yo'q"),
    ("IV bosqich:", "Sof simulyakr — o'z-o'zidan amal qiluvchi belgi"),
]

y = 1.55
for label, desc in stages_text:
    add_text_block(slide, Inches(9.2), Inches(y), Inches(3.6), Inches(0.35),
                   label, font_size=12, bold=True, color=ACCENT_RED)
    add_text_block(slide, Inches(9.2), Inches(y + 0.35), Inches(3.6), Inches(0.6),
                   desc, font_size=11, color=TEXT_DARK)
    y += 1.0

# Bottom note
note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.65))
note.fill.solid()
note.fill.fore_color.rgb = ACCENT_GOLD
note.line.fill.background()
add_text_block(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
               "MUHIM XULOSA: Hozirgi virtual reallik (VR) — bu IV bosqichdagi sof simulyakrning texnologik mujassamlanishidir. "
               "U asl haqiqatga emas, o'zining ichki kodiga asoslanadi.",
               font_size=12, bold=True, color=DARK_BLUE)


# ===================== SLIDE 8: VIRTUAL REALLIK ONTOLOGIK MAQOMI =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "II BOB. VIRTUAL REALLIKNING ONTOLOGIK MAQOMI", slide_num=8)
add_footer(slide)

# Insert conceptual scheme
slide.shapes.add_picture('charts/chart6_conceptual_scheme.png',
                          Inches(0.4), Inches(1.05),
                          width=Inches(12.5), height=Inches(4.2))

# Bottom analysis
add_text_block(slide, Inches(0.4), Inches(5.35), Inches(12.5), Inches(0.4),
               "VIRTUAL REALLIKNING ONTOLOGIK PARAMETRLARI:",
               font_size=14, bold=True, color=DARK_BLUE)

params = [
    ("IMMATERIALLIK", "Moddiy substansiyaga ega emas,\nlekin haqiqiy ta'sir ko'rsatadi", TEAL),
    ("INTERAKTIVLIK", "Foydalanuvchi bilan o'zaro\nta'sirga kirisha oladi", ACCENT_RED),
    ("IMMERSIVLIK", "Inson ongini to'liq\no'ziga jalb qiladi", ACCENT_GOLD),
    ("AVTONOMIYA", "Asl haqiqatga muhtoj emas,\no'zi reallik yaratadi", DARK_BLUE),
]

for i, (name, desc, color) in enumerate(params):
    left = Inches(0.4 + i * 3.18)
    top = Inches(5.85)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(3.0), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)

    add_text_block(slide, left, top + Inches(0.1), Inches(3.0), Inches(0.4),
                   name, font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_block(slide, left + Inches(0.15), top + Inches(0.5),
                   Inches(2.7), Inches(0.7),
                   desc, font_size=10.5, color=TEXT_DARK, align=PP_ALIGN.CENTER)


# ===================== SLIDE 9: AXLOQIY IMPLIKATSIYALAR (CHART) =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "VIRTUAL REALLIKNING AXLOQIY IMPLIKATSIYALARI", slide_num=9)
add_footer(slide)

# Chart
slide.shapes.add_picture('charts/chart3_ethical_problems.png',
                          Inches(0.3), Inches(1.0),
                          width=Inches(7.0), height=Inches(5.5))

# Right side - dilemmas
add_text_block(slide, Inches(7.6), Inches(1.05), Inches(5.4), Inches(0.4),
               "ASOSIY AXLOQIY DILEMMALAR",
               font_size=15, bold=True, color=DARK_BLUE)

dilemmas = [
    ("Avatar paradoksi:",
     "Virtual avatar orqali sodir etilgan harakat — kimning javobgarligi?"),
    ("Identitet krizisi:",
     "Bir nechta virtual «men» — hammasi haqiqiymi?"),
    ("Mas'uliyat tarqalishi:",
     "Algoritmik qaror — kim mas'ul: dasturchimi, mashinami, foydalanuvchimi?"),
    ("Empatiya yo'qolishi:",
     "Virtual qurbon — chinakam qurbonmi?"),
    ("Axloqiy relativizm:",
     "Har bir virtual olamning o'z axloqi bo'lsa, umumbashariy axloq qoladimi?"),
]

y = 1.55
for title, desc in dilemmas:
    add_text_block(slide, Inches(7.7), Inches(y), Inches(5.3), Inches(0.4),
                   "■  " + title, font_size=12, bold=True, color=ACCENT_RED)
    add_text_block(slide, Inches(8.0), Inches(y + 0.3), Inches(5.0), Inches(0.7),
                   desc, font_size=11, italic=True, color=TEXT_DARK)
    y += 1.0

# Bottom highlight
hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5))
hl.fill.solid()
hl.fill.fore_color.rgb = ACCENT_RED
hl.line.fill.background()
add_text_block(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
               "Bodriyarning prognozi: virtual axloq — bu axloqning yo'qolishi emas, balki uning «giperreal kodlash»idir",
               font_size=12.5, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# ===================== SLIDE 10: VIRTUAL ZO'RAVONLIK (CHART) =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "2.2. RAQAMLI GIPERREALLIK VA VIRTUAL ZO'RAVONLIK", slide_num=10)
add_footer(slide)

# VR users dynamics chart
slide.shapes.add_picture('charts/chart2_vr_dynamics.png',
                          Inches(0.3), Inches(1.05),
                          width=Inches(8.0), height=Inches(4.4))

# Right side - virtual violence types
add_text_block(slide, Inches(8.5), Inches(1.05), Inches(4.6), Inches(0.4),
               "VIRTUAL ZO'RAVONLIK TURLARI",
               font_size=14, bold=True, color=ACCENT_RED)

types = [
    ("1. Simulyativ", "VR o'yinlardagi zo'ravonlik\n(GTA, FPS o'yinlari)"),
    ("2. Immersiv", "VR-da to'liq sho'ng'igan\nholda kechiriladigan tajriba"),
    ("3. Algoritmik", "AI tomonidan amalga\noshirilgan psixologik bosim"),
    ("4. Onto-axloqiy", "Identitet va mavjudlikka\nqilingan tahdid"),
]

y = 1.55
for title, desc in types:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(8.5), Inches(y), Inches(4.5), Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = ACCENT_RED
    box.line.width = Pt(1.5)

    add_text_block(slide, Inches(8.65), Inches(y + 0.05), Inches(4.3), Inches(0.3),
                   title, font_size=12, bold=True, color=ACCENT_RED)
    add_text_block(slide, Inches(8.65), Inches(y + 0.32), Inches(4.3), Inches(0.55),
                   desc, font_size=10.5, color=TEXT_DARK)
    y += 0.95

# Bottom analytical block
add_decorated_box(slide, Inches(0.4), Inches(5.55), Inches(12.5), Inches(1.55),
                  "BODRIYAR PARADOKSI: Virtualda yo'q narsa, real ta'sir ko'rsatadi",
                  ["Empirik dalillar: VR-zo'ravonlikni ko'rgan inson real PTSD belgilari namoyish qiladi (Stanford Univ., 2023)",
                   "Ontologik nuance: virtual zo'ravonlikning «chegarasi» yo'q — qurbon ham, jallod ham simulyakrdir",
                   "Axloqiy paradoks: jismoniy zarar yo'q, lekin onglar transformatsiyasi va begonalashuv yuzaga keladi",
                   "Yangilik: 2024-yilda BMT Inson huquqlari kengashi virtual zo'ravonlikni rasman tan oldi"],
                  title_color=ACCENT_RED)


# ===================== SLIDE 11: SIMULYATSIYA - IJTIMOIY-SIYOSIY =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "III BOB. SIMULYATSIYANING IJTIMOIY-SIYOSIY ROLI", slide_num=11)
add_footer(slide)

# Chart on left
slide.shapes.add_picture('charts/chart4_simulation_spheres.png',
                          Inches(0.3), Inches(1.0),
                          width=Inches(8.0), height=Inches(4.5))

# Right column
add_text_block(slide, Inches(8.5), Inches(1.05), Inches(4.6), Inches(0.4),
               "SIMULYATSIYA NAMOYON BO'LISH SOHALARI",
               font_size=13, bold=True, color=DARK_BLUE)

spheres_list = [
    ("Saylov:", "Imidj > dastur, brending > g'oya"),
    ("Media:", "Yangilik faktdan ko'ra dramaturgiya"),
    ("Iste'mol:", "Mahsulot emas — orzu sotiladi"),
    ("Geosiyosat:", "«Xalifalik urushi» (Bodriyar, 1991)"),
    ("Ijtimoiy tarmoq:", "Profil > shaxs, lijklar > qadr"),
    ("Ta'lim:", "Ko'rsatkichlar > bilim mazmuni"),
]

y = 1.55
for title, desc in spheres_list:
    add_text_block(slide, Inches(8.55), Inches(y), Inches(4.5), Inches(0.32),
                   "▸ " + title, font_size=12, bold=True, color=ACCENT_RED)
    add_text_block(slide, Inches(8.85), Inches(y + 0.3), Inches(4.2), Inches(0.35),
                   desc, font_size=11, italic=True, color=TEXT_DARK)
    y += 0.65

# Bottom - Bodriyar's famous case
case_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.5))
case_box.fill.solid()
case_box.fill.fore_color.rgb = DARK_BLUE
case_box.line.color.rgb = ACCENT_GOLD
case_box.line.width = Pt(2)

add_text_block(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.4),
               'KEYS: «KO\'RFAZ URUSHI BO\'LIB O\'TMADI» (BODRIYAR, 1991)',
               font_size=14, bold=True, color=ACCENT_GOLD)
add_text_block(slide, Inches(0.6), Inches(6.1), Inches(12.1), Inches(1.0),
               "Bodriyarning mashhur essesida 1991-yilgi Ko'rfaz urushi haqiqiy urush emas, balki CNN tomonidan teleekranda "
               "yaratilgan simulyakr ekanligi ta'kidlanadi. Tomoshabin uchun urush — bu video o'yin estetikasi, ko'r raketalarning "
               "yorug' chizig'i, lekin asl o'lim, qon va azob — ekrandan tashqarida qoladi. Bu — simulyatsiyaning eng yorqin "
               "siyosiy misoli.",
               font_size=11.5, color=WHITE)


# ===================== SLIDE 12: OMMAVIY KOMMUNIKATSIYA =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "3.1. OMMAVIY KOMMUNIKATSIYADA SIMULYATSIYA", slide_num=12)
add_footer(slide)

# Top concept
add_text_block(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.5),
               "Bodriyar bo'yicha media — bu kommunikatsiya emas, balki kommunikatsiyaning simulyatsiyasidir",
               font_size=14, bold=True, italic=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Three columns
columns = [
    ("AN'ANAVIY MEDIA",
     ["Bir tomonlama ma'lumot",
      "Faktga asoslangan",
      "Manba ko'rsatiladi",
      "Saralash mavjud",
      "Vaqt sekin"],
     "(matn — gazeta — TV)",
     TEAL),
    ("RAQAMLI MEDIA",
     ["Interaktiv platforma",
      "Faktdan ko'ra hissiyot",
      "Manbalar aralash",
      "Saralash algoritmik",
      "Real-vaqt rejimi"],
     "(internet — ijtimoiy tarmoq)",
     ACCENT_GOLD),
    ("GIPERREAL MEDIA",
     ["Foydalanuvchi — kontent yaratuvchisi",
      "Tarkibsiz spektakl",
      "Manba — ALGORITM",
      "Filtrlangan reallik",
      "Vaqt yo'qoladi"],
     "(TikTok, AI, deepfake)",
     ACCENT_RED),
]

for i, (title, items, sub, color) in enumerate(columns):
    left = Inches(0.4 + i * 4.25)
    top = Inches(1.7)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, Inches(4.05), Inches(4.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, top, Inches(4.05), Inches(0.7))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    add_text_block(slide, left, top + Inches(0.1), Inches(4.05), Inches(0.4),
                   title, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_block(slide, left, top + Inches(0.4), Inches(4.05), Inches(0.3),
                   sub, font_size=10, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Items
    for j, item in enumerate(items):
        add_text_block(slide, left + Inches(0.2), top + Inches(0.85 + j * 0.6),
                       Inches(3.65), Inches(0.55),
                       "✓  " + item, font_size=11.5, color=TEXT_DARK)

# Bottom conclusion
hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.85))
hl.fill.solid()
hl.fill.fore_color.rgb = DARK_BLUE
hl.line.fill.background()
add_text_block(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.4),
               "BODRIYARNING ASOSIY TEZISI:",
               font_size=12, bold=True, color=ACCENT_GOLD)
add_text_block(slide, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.4),
               "«Axborot ko'paygan sari, ma'no kamayadi» — informatsion o'sish ma'no inflatsiyasini keltirib chiqaradi.",
               font_size=12, italic=True, color=WHITE)


# ===================== SLIDE 13: BODRIYAR TANQIDI =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "3.2. BODRIYAR G'OYALARINING TANQIDIY TAHLILI", slide_num=13)
add_footer(slide)

# Two columns: positive and critical
# Left - strengths
strength_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.4), Inches(1.05), Inches(6.2), Inches(2.8))
strength_box.fill.solid()
strength_box.fill.fore_color.rgb = WHITE
strength_box.line.color.rgb = TEAL
strength_box.line.width = Pt(2)

header1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.4), Inches(1.05), Inches(6.2), Inches(0.5))
header1.fill.solid()
header1.fill.fore_color.rgb = TEAL
header1.line.fill.background()
add_text_block(slide, Inches(0.4), Inches(1.15), Inches(6.2), Inches(0.4),
               "✚  NAZARIYANING KUCHLI TOMONLARI",
               font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

strengths = [
    "Raqamli asrning ontologik tahlili uchun tayanch nazariya",
    "Iste'mol madaniyatining tanqidiy talqini",
    "Media-tahlilning sivilizatsion miqyosi",
    "Postmodern hodisalarni bashorat qilish qudrati",
    "Belgilar tizimining falsafiy konseptualizatsiyasi",
]

for i, s in enumerate(strengths):
    add_text_block(slide, Inches(0.6), Inches(1.7 + i * 0.4), Inches(5.9), Inches(0.4),
                   "✓  " + s, font_size=11.5, color=TEXT_DARK)

# Right - weaknesses
weak_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(6.7), Inches(1.05), Inches(6.2), Inches(2.8))
weak_box.fill.solid()
weak_box.fill.fore_color.rgb = WHITE
weak_box.line.color.rgb = ACCENT_RED
weak_box.line.width = Pt(2)

header2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(6.7), Inches(1.05), Inches(6.2), Inches(0.5))
header2.fill.solid()
header2.fill.fore_color.rgb = ACCENT_RED
header2.line.fill.background()
add_text_block(slide, Inches(6.7), Inches(1.15), Inches(6.2), Inches(0.4),
               "✕  TANQIDIY ZAIFLIKLAR",
               font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

weaknesses = [
    "Texno-pessimizm — alternativalar ko'rsatilmagan",
    "Empirik isbotlar yetishmaydi (essecistik uslub)",
    "Universalizm — Sharq madaniyati hisobga olinmagan",
    "Inson agentligi rolini kamaytirish",
    "Determinizm — texnologiya hal qiluvchi omil sifatida",
]

for i, w in enumerate(weaknesses):
    add_text_block(slide, Inches(6.9), Inches(1.7 + i * 0.4), Inches(5.9), Inches(0.4),
                   "✗  " + w, font_size=11.5, color=TEXT_DARK)

# Critics list
add_text_block(slide, Inches(0.4), Inches(4.05), Inches(12.5), Inches(0.4),
               "ASOSIY TANQIDCHILAR VA MUQOBIL POZITSIYALAR:",
               font_size=14, bold=True, color=DARK_BLUE)

critics = [
    ("D. Kelner", "Tanqidiy nazariya pozitsiyasidan: Bodriyar siyosatni e'tiborsiz qoldirgan", ACCENT_GOLD),
    ("M. Gane", "Bodriyar nazariyasi falsafa emas, badiiy provokatsiyadir", RGBColor(0x8E, 0x44, 0xAD)),
    ("S. Best", "Postmodern apoteoz — yangi diskursiv strategiyalar kerak", TEAL),
    ("M. Castells", "Tarmoq jamiyati nazariyasi — muqobil va konstruktivroq", DARK_BLUE),
]

for i, (name, view, color) in enumerate(critics):
    row = i // 2
    col = i % 2
    left = Inches(0.4 + col * 6.3)
    top = Inches(4.55 + row * 1.25)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(6.2), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    # Name badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left + Inches(0.15), top + Inches(0.15),
                                     Inches(1.4), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_text_block(slide, left + Inches(0.15), top + Inches(0.22),
                   Inches(1.4), Inches(0.4),
                   name, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_block(slide, left + Inches(1.65), top + Inches(0.18),
                   Inches(4.45), Inches(0.85),
                   view, font_size=11.5, italic=True, color=TEXT_DARK)


# ===================== SLIDE 14: INSONIY BEGONALASHUV (CHART) =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_header_bar(slide, "VIRTUAL REALLIK VA INSONIY BEGONALASHUV", slide_num=14)
add_footer(slide)

# Chart on left
slide.shapes.add_picture('charts/chart5_alienation.png',
                          Inches(0.2), Inches(1.0),
                          width=Inches(7.5), height=Inches(5.4))

# Right column - typology
add_text_block(slide, Inches(7.9), Inches(1.05), Inches(5.2), Inches(0.4),
               "RAQAMLI BEGONALASHUV TIPOLOGIYASI",
               font_size=13, bold=True, color=DARK_BLUE)

types = [
    ("Onto-begonalashuv:", "O'z mavjudligidan uzilish.\nAvatar — chinakam «men»dan ustun"),
    ("Aksiologik begonalashuv:", "Qadriyatlar orientatsiyasini\nyo'qotish, axloqiy relyativizm"),
    ("Kommunikativ begonalashuv:", "Boshqalarsiz «do'stlik» —\nemoji va laykning fonologiyasi"),
    ("Temporal begonalashuv:", "Chinakam vaqtdan tushish —\ncheksiz «hozir» (real-time)"),
    ("Tana begonalashuvi:", "Jismoniy tanani inkor etish,\nsof raqamli ekzistensiya orzu"),
]

y = 1.55
for title, desc in types:
    add_text_block(slide, Inches(8.0), Inches(y), Inches(5.0), Inches(0.32),
                   "■  " + title, font_size=12, bold=True, color=ACCENT_RED)
    add_text_block(slide, Inches(8.3), Inches(y + 0.3), Inches(4.7), Inches(0.65),
                   desc, font_size=10.5, italic=True, color=TEXT_DARK)
    y += 1.0

# Bottom - solution direction
sol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.55))
sol.fill.solid()
sol.fill.fore_color.rgb = TEAL
sol.line.fill.background()
add_text_block(slide, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.45),
               "MUQOBIL YO'L: «Onglilik bilan virtuallashuv» — tanqidiy yondashuv, raqamli detoks va onto-axloqiy ta'lim",
               font_size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ===================== SLIDE 15: XULOSA =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_header_bar(slide, "XULOSA: TADQIQOTNING ASOSIY NATIJALARI", slide_num=15)

# Footer
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), prs.slide_width, Inches(0.35))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GOLD
bar.line.fill.background()
add_text_block(slide, Inches(0.4), Inches(7.18), Inches(12.5), Inches(0.3),
               "E'TIBORINGIZ UCHUN RAHMAT!  |  THANK YOU FOR YOUR ATTENTION!",
               font_size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Main conclusions box
add_text_block(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.4),
               "ASOSIY ILMIY NATIJALAR:",
               font_size=16, bold=True, color=ACCENT_GOLD)

conclusions = [
    ("1.", "Bodriyar simulyakr nazariyasi XXI asr raqamli muhitining ontologik strukturasini tushunish uchun "
           "yetakchi falsafiy paradigma sifatida tasdiqlandi."),
    ("2.", "Virtual reallik IV-bosqich simulyakri sifatida ontologik avtonomiyaga ega bo'lib, asl haqiqatga "
           "muhtoj bo'lmagan holda mustaqil reallik yaratadi."),
    ("3.", "Raqamli giperreallik yangi tipdagi axloqiy dilemmalarni keltirib chiqaradi: avatar paradoksi, "
           "algoritmik mas'uliyat, virtual zo'ravonlik va onto-aksiologik begonalashuv."),
    ("4.", "Bodriyar nazariyasi texno-pessimistik yo'nalishga ega bo'lsa-da, tanqidiy qayta ko'rib chiqilganda "
           "Sharq aksiologiyasi bilan boyitiladi va konstruktiv yo'nalish oladi."),
    ("5.", "Tadqiqot natijasida raqamli muhitda axloqiy ta'lim uchun «onto-aksiologik" ' grammatika» ' "modeli ishlab chiqildi."),
]

y = 1.55
for num, text in conclusions:
    # Number box
    box = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.5), Inches(y), Inches(0.5), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_GOLD
    box.line.fill.background()
    add_text_block(slide, Inches(0.5), Inches(y + 0.05), Inches(0.5), Inches(0.4),
                   num, font_size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Text
    add_text_block(slide, Inches(1.15), Inches(y + 0.05), Inches(11.7), Inches(0.85),
                   text, font_size=12, color=WHITE)
    y += 1.0

# Final quote
quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.5))
quote.fill.solid()
quote.fill.fore_color.rgb = ACCENT_RED
quote.line.color.rgb = ACCENT_GOLD
quote.line.width = Pt(1.5)
add_text_block(slide, Inches(1.7), Inches(6.6), Inches(9.9), Inches(0.4),
               '"Haqiqat shu darajada ko\'rinadiki, u endi ko\'rinmas." — J. Bodriyar',
               font_size=12.5, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============ SAVE ============
output_path = "Bodriyar_Virtual_Reallik_Taqdimot.pptx"
prs.save(output_path)
print(f"Taqdimot yaratildi: {output_path}")
print(f"Slaydlar soni: {len(prs.slides)}")

# File size
import os
size_kb = os.path.getsize(output_path) / 1024
print(f"Hajmi: {size_kb:.1f} KB")
